"""
build_clip_slide.py
-------------------
Single-slide PPTX — CLIP pipeline diagram matching the screenshot style.
Run:  python build_clip_slide.py
Out:  CLIP_slide.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

W = Inches(13.33)
H = Inches(7.5)

NAVY  = RGBColor(0x1A, 0x2E, 0x4A)
BLUE  = RGBColor(0x4A, 0x6F, 0x9A)   # slate-blue matching screenshot
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x22, 0x22, 0x22)
LGREY = RGBColor(0xF0, 0xF2, 0xF5)
VBAR  = RGBColor(0xB0, 0xC8, 0xE0)   # vector bar colour

TF = "Calibri"

def prs_new():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=LGREY, line_c=None, line_w=0):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_c:
        s.line.color.rgb = line_c; s.line.width = line_w
    else:
        s.line.width = 0
    return s

def lbl(slide, text, l, t, w, h, size=Pt(16), bold=False,
        color=DARK, align=PP_ALIGN.CENTER, italic=False):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = text
    r.font.name = TF; r.font.size = size
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    p.alignment = align
    return b

def arrow(slide, l, t, length, thickness=Pt(3)):
    """Simple horizontal arrow line."""
    ln = slide.shapes.add_shape(1, l, t, length, thickness)
    ln.fill.solid(); ln.fill.fore_color.rgb = DARK
    ln.line.width = 0

def main():
    prs = prs_new()
    slide = blank(prs)

    # White background
    rect(slide, 0, 0, W, H, fill=WHITE)

    # Title bar
    rect(slide, 0, 0, W, Inches(1.1), fill=NAVY)
    lbl(slide, "CLIP Visual Verification — How It Works",
        Inches(0.5), Inches(0.15), W - Inches(1), Inches(0.85),
        size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # ── Layout measurements ───────────────────────────────────────────────────
    # 5 nodes: [img] → [CLIP] → [vector] → [probe] → [output]
    # Centered vertically in the slide

    cy   = Inches(3.5)       # center Y of diagram
    box_h = Inches(1.6)
    box_w = Inches(2.0)
    vbar_w = Inches(0.7)
    gap  = Inches(0.6)       # gap between elements (for arrow)

    total_w = box_w + gap + box_w + gap + vbar_w + gap + box_w + gap + box_w
    start_x = (W - total_w) / 2

    x = start_x
    top = cy - box_h / 2

    # ── 1. Gripper image placeholder ─────────────────────────────────────────
    img_box_path = os.path.join(
        os.path.dirname(__file__),
        "C:\\Users\\myuwa\\.gemini\\antigravity\\brain\\fbdf0c36-6916-4be2-97ac-c9e5321e04d1\\clip_flow_diagram_hires_1773922628240.png"
    )
    # Draw a grey placeholder box for the gripper image
    rect(slide, x, top, box_w, box_h,
         fill=RGBColor(0xE0,0xE5,0xEC),
         line_c=RGBColor(0xAA,0xBB,0xCC), line_w=Pt(1.5))
    lbl(slide, "[Gripper\nImage]", x, top + Inches(0.45), box_w, Inches(0.8),
        size=Pt(17), bold=True, color=BLUE)
    lbl(slide, "Gripper Image Crop",
        x, top + box_h + Inches(0.1), box_w, Inches(0.45),
        size=Pt(13), color=DARK)

    # Arrow
    x += box_w
    arrow(slide, x + Inches(0.05), cy - Pt(1.5), gap - Inches(0.12))
    lbl(slide, "→", x, cy - Inches(0.32), gap, Inches(0.65),
        size=Pt(30), bold=True, color=DARK)
    x += gap

    # ── 2. CLIP Model box ────────────────────────────────────────────────────
    rect(slide, x, top, box_w, box_h, fill=BLUE)
    lbl(slide, "CLIP\nModel", x, top + Inches(0.3), box_w, Inches(0.9),
        size=Pt(22), bold=True, color=WHITE)
    lbl(slide, "ViT-B/32", x, top + Inches(1.15), box_w, Inches(0.35),
        size=Pt(13), color=RGBColor(0xD0,0xE4,0xF4))

    x += box_w
    lbl(slide, "→", x, cy - Inches(0.32), gap, Inches(0.65),
        size=Pt(30), bold=True, color=DARK)
    x += gap

    # ── 3. 512-dim Vector bar ────────────────────────────────────────────────
    bar_top = top + Inches(0.15)
    bar_h_total = box_h - Inches(0.3)
    n_bars = 16
    seg_h = bar_h_total / n_bars
    colors = [
        RGBColor(0x1A,0x4A,0x7A),RGBColor(0x2D,0x6A,0xA0),
        RGBColor(0x4A,0x8F,0xC0),RGBColor(0x6A,0xAA,0xCF),
        RGBColor(0xA0,0xD0,0xEA),RGBColor(0xC8,0xE8,0xF5),
        RGBColor(0xE8,0xF4,0xFA),RGBColor(0xF5,0xF8,0xFA),
        RGBColor(0xFA,0xF0,0xD5),RGBColor(0xF5,0xD8,0x90),
        RGBColor(0xF0,0xB8,0x50),RGBColor(0xE0,0x90,0x30),
        RGBColor(0xC0,0x60,0x20),RGBColor(0x9A,0x40,0x18),
        RGBColor(0x70,0x28,0x10),RGBColor(0x50,0x18,0x08),
    ]
    for i in range(n_bars):
        rect(slide, x, bar_top + i * seg_h, vbar_w, seg_h,
             fill=colors[i % len(colors)])
    rect(slide, x, bar_top, vbar_w, bar_h_total,
         fill=RGBColor(0,0,0), line_c=RGBColor(0x80,0x80,0x80), line_w=Pt(0.5))
    rect(slide, x, bar_top, vbar_w, bar_h_total,
         fill=RGBColor(0,0,0,))
    # Redraw properly — segments then border
    for i in range(n_bars):
        s = slide.shapes.add_shape(1, x, bar_top + i * seg_h, vbar_w, seg_h)
        s.fill.solid(); s.fill.fore_color.rgb = colors[i % len(colors)]
        s.line.width = 0
    # Border
    bdr = slide.shapes.add_shape(1, x, bar_top, vbar_w, bar_h_total)
    bdr.fill.background()
    bdr.line.color.rgb = RGBColor(0x80,0x80,0x80); bdr.line.width = Pt(0.8)

    lbl(slide, "[512-dim\nVector]",
        x - Inches(0.15), top + box_h + Inches(0.1), vbar_w + Inches(0.3), Inches(0.55),
        size=Pt(13), color=DARK)

    x += vbar_w
    lbl(slide, "→", x, cy - Inches(0.32), gap, Inches(0.65),
        size=Pt(30), bold=True, color=DARK)
    x += gap

    # ── 4. Linear Probe box ──────────────────────────────────────────────────
    rect(slide, x, top, box_w, box_h, fill=BLUE)
    lbl(slide, "Linear\nProbe", x, top + Inches(0.3), box_w, Inches(0.9),
        size=Pt(22), bold=True, color=WHITE)
    lbl(slide, "Logistic\nRegression", x, top + Inches(1.1), box_w, Inches(0.42),
        size=Pt(12), color=RGBColor(0xD0,0xE4,0xF4))

    x += box_w
    lbl(slide, "→", x, cy - Inches(0.32), gap, Inches(0.65),
        size=Pt(30), bold=True, color=DARK)
    x += gap

    # ── 5. Output boxes ───────────────────────────────────────────────────────
    out_h = Inches(0.7)
    out_gap = Inches(0.25)
    out_top_y = cy - out_h - out_gap / 2

    # "Holding" ✓
    rect(slide, x, out_top_y, box_w, out_h, fill=RGBColor(0x1E,0x7E,0x45))
    lbl(slide, '"Holding"  ✓', x, out_top_y + Inches(0.1), box_w, out_h - Inches(0.1),
        size=Pt(17), bold=True, color=WHITE)

    # "Empty" ✗
    rect(slide, x, out_top_y + out_h + out_gap, box_w, out_h,
         fill=RGBColor(0xC8,0x35,0x35))
    lbl(slide, '"Empty"  ✗', x, out_top_y + out_h + out_gap + Inches(0.1),
        box_w, out_h - Inches(0.1),
        size=Pt(17), bold=True, color=WHITE)

    # Caption
    lbl(slide, "or", x, cy - Inches(0.18), box_w, Inches(0.35),
        size=Pt(14), color=DARK, bold=True)

    # ── Footer note ───────────────────────────────────────────────────────────
    lbl(slide,
        "CLIP extracts a 512-dimensional image embedding → Logistic Regression (linear probe) trained on ~600 gripper images classifies it.",
        Inches(0.5), Inches(6.6), W - Inches(1), Inches(0.6),
        size=Pt(13), color=RGBColor(0x55,0x55,0x55), italic=True,
        align=PP_ALIGN.CENTER)

    prs.save("CLIP_slide.pptx")
    print("\n  Saved: CLIP_slide.pptx\n")

if __name__ == "__main__":
    main()
