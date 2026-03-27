"""
build_clip_diagram.py
---------------------
Single slide PPTX — exact match to the screenshot:
  [Gripper photo] → [CLIP Model] → [|||] → [Linear Probe] → "Holding" or "Empty"

Run:  python build_clip_diagram.py
Out:  CLIP_diagram.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

W = Inches(13.33)
H = Inches(7.5)

# Exact colours from screenshot
STEEL   = RGBColor(0x44, 0x72, 0xC4)   # the steel blue of the boxes
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x26, 0x26, 0x26)
LGREY   = RGBColor(0xF2, 0xF2, 0xF2)
ARROW_C = RGBColor(0x40, 0x40, 0x40)
TF = "Calibri"

def prs_new():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, l, t, w, h, fill=WHITE, line_c=None, line_w=0):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_c:
        s.line.color.rgb = line_c; s.line.width = line_w
    else:
        s.line.width = 0
    return s

def add_pentagon(slide, l, t, w, h, fill=STEEL):
    """Pentagon shape (right-pointing) to match screenshot chevron style."""
    # MSO auto shape: Pentagon = 56
    from pptx.util import Emu
    sp = slide.shapes.add_shape(56, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.width = 0
    return sp

def add_txt(slide, text, l, t, w, h, size=Pt(16), bold=False,
            color=DARK, align=PP_ALIGN.CENTER, italic=False):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = text
    r.font.name = TF; r.font.size = size
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    p.alignment = align
    return b

def multiline(slide, lines, l, t, w, h, size=Pt(16), bold=False,
              color=DARK, align=PP_ALIGN.CENTER):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.name = TF; r.font.size = size
        r.font.bold = bold; r.font.color.rgb = color
        p.alignment = align
    return b

def add_arrow(slide, l, t, w, h=Inches(0.25)):
    """Right-pointing block arrow."""
    s = slide.shapes.add_shape(13, l, t, w, h)   # 13 = RIGHT_ARROW
    s.fill.solid(); s.fill.fore_color.rgb = ARROW_C
    s.line.width = 0
    return s


def main():
    prs = prs_new()
    slide = blank(prs)

    # Pure white background
    add_rect(slide, 0, 0, W, H, fill=WHITE)

    # ── Layout constants ───────────────────────────────────────────────────────
    # Centre everything vertically
    CY = Inches(3.55)          # vertical centre
    BOX_H = Inches(1.5)        # height of blue boxes
    BOX_W = Inches(2.0)        # width of blue boxes
    IMG_W = Inches(1.8)        # gripper image placeholder
    IMG_H = Inches(1.5)
    BAR_W = Inches(0.55)       # width of vector bar
    BAR_H = Inches(1.5)
    ARR_W = Inches(0.7)        # arrow width
    ARR_H = Inches(0.22)
    GAP   = Inches(0.12)       # gap between arrow and box
    OUT_W = Inches(1.5)        # output text box
    LBL_H = Inches(0.45)       # label below height

    # Total width of diagram
    total = IMG_W + ARR_W + BOX_W + ARR_W + BAR_W + ARR_W + BOX_W + ARR_W + OUT_W
    sx = (W - total) / 2       # start x
    ty = CY - BOX_H / 2        # top y

    x = sx

    # ── 1. Gripper image placeholder ─────────────────────────────────────────
    # Light grey box with a small icon-style robot arm drawing
    add_rect(slide, x, ty, IMG_W, IMG_H,
             fill=RGBColor(0xD5, 0xD5, 0xD5),
             line_c=RGBColor(0x99, 0x99, 0x99), line_w=Pt(1))

    # Simple gripper icon using shapes (two rectangles = finger + palm)
    ix = x + Inches(0.4); iy = ty + Inches(0.15)
    # Palm
    add_rect(slide, ix + Inches(0.2), iy + Inches(0.25),
             Inches(0.65), Inches(0.55), fill=RGBColor(0x60, 0x60, 0x60))
    # Left finger
    add_rect(slide, ix + Inches(0.1), iy + Inches(0.8),
             Inches(0.22), Inches(0.42), fill=RGBColor(0x50, 0x50, 0x50))
    # Right finger
    add_rect(slide, ix + Inches(0.72), iy + Inches(0.8),
             Inches(0.22), Inches(0.42), fill=RGBColor(0x50, 0x50, 0x50))
    # Object between fingers (purple like screenshot)
    add_rect(slide, ix + Inches(0.32), iy + Inches(0.82),
             Inches(0.4), Inches(0.35), fill=RGBColor(0x80, 0x40, 0xA0))

    # Label below
    add_txt(slide, "[Gripper\nImage]",
            x, ty + IMG_H + Inches(0.08), IMG_W, LBL_H + Inches(0.15),
            size=Pt(14), color=DARK, bold=False)

    x += IMG_W

    # ── Arrow 1 ───────────────────────────────────────────────────────────────
    add_arrow(slide, x + GAP, CY - ARR_H / 2, ARR_W - GAP * 2, ARR_H)
    x += ARR_W

    # ── 2. CLIP Model (pentagon) ──────────────────────────────────────────────
    add_pentagon(slide, x, ty, BOX_W, BOX_H, fill=STEEL)
    multiline(slide, ["CLIP", "Model"], x, ty + Inches(0.3), BOX_W, Inches(0.85),
              size=Pt(18), bold=True, color=WHITE)
    x += BOX_W

    # ── Arrow 2 ───────────────────────────────────────────────────────────────
    add_arrow(slide, x + GAP, CY - ARR_H / 2, ARR_W - GAP * 2, ARR_H)
    x += ARR_W

    # ── 3. 512-dim Vector bars ────────────────────────────────────────────────
    bar_top = ty + Inches(0.05)
    n = 14
    seg = BAR_H / n
    # Gradient from dark blue to light blue (like screenshot grayscale bars)
    for i in range(n):
        shade = int(80 + (i / n) * 140)
        col = RGBColor(shade, shade + 10, min(255, shade + 40))
        s = slide.shapes.add_shape(1, x, bar_top + i * seg,
                                   BAR_W, seg + Pt(0.5))
        s.fill.solid(); s.fill.fore_color.rgb = col
        s.line.width = 0
    # Thin border over bars
    bdr = slide.shapes.add_shape(1, x, bar_top, BAR_W, BAR_H)
    bdr.fill.background()
    bdr.line.color.rgb = RGBColor(0x88, 0x88, 0x88)
    bdr.line.width = Pt(0.75)

    # Label below
    multiline(slide, ["[512-dim", "Vector]"], x - Inches(0.1),
              ty + BAR_H + Inches(0.08), BAR_W + Inches(0.2), LBL_H + Inches(0.15),
              size=Pt(14), color=DARK)
    x += BAR_W

    # ── Arrow 3 ───────────────────────────────────────────────────────────────
    add_arrow(slide, x + GAP, CY - ARR_H / 2, ARR_W - GAP * 2, ARR_H)
    x += ARR_W

    # ── 4. Linear Probe (pentagon) ────────────────────────────────────────────
    add_pentagon(slide, x, ty, BOX_W, BOX_H, fill=STEEL)
    multiline(slide, ["Linear", "Probe"], x, ty + Inches(0.3), BOX_W, Inches(0.85),
              size=Pt(18), bold=True, color=WHITE)
    x += BOX_W

    # ── Arrow 4 ───────────────────────────────────────────────────────────────
    add_arrow(slide, x + GAP, CY - ARR_H / 2, ARR_W - GAP * 2, ARR_H)
    x += ARR_W

    # ── 5. Output text ────────────────────────────────────────────────────────
    multiline(slide,
              ['"Holding"', 'or', '"Empty"'],
              x, CY - Inches(0.7), OUT_W, Inches(1.4),
              size=Pt(17), bold=True, color=DARK, align=PP_ALIGN.LEFT)

    prs.save("CLIP_diagram.pptx")
    print("\n  Saved: CLIP_diagram.pptx\n")

if __name__ == "__main__":
    main()
