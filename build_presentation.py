"""
build_presentation.py
---------------------
Generates the FYP presentation PPTX.
Run from any directory:  python build_presentation.py
Output: FYP_Presentation.pptx  (same folder as this script)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x2E, 0x4A)   # slide title bar, headings
MID_BL  = RGBColor(0x2D, 0x5F, 0x8A)   # accent rules, table header
LIGHT   = RGBColor(0xF4, 0xF6, 0xF9)   # slide background tint (very light grey)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1C, 0x1C, 0x1C)   # body text
ACCENT  = RGBColor(0x2D, 0x8A, 0x4E)   # green tick / highlight
AMBER   = RGBColor(0xD4, 0x7D, 0x00)   # warning / WIP

TITLE_FONT   = "Calibri"
BODY_FONT    = "Calibri"

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def bg_rect(slide, color=LIGHT):
    """Fill entire slide with a solid background colour."""
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, H
    )
    shape.line.color.rgb = color
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.width = 0


def title_bar(slide, text, bar_h=Inches(1.25)):
    """Dark navy bar at top with white title text."""
    bar = slide.shapes.add_shape(1, 0, 0, W, bar_h)
    bar.fill.solid();  bar.fill.fore_color.rgb = NAVY
    bar.line.width = 0

    tf = bar.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = text
    r.font.name  = TITLE_FONT
    r.font.size  = Pt(32)
    r.font.bold  = True
    r.font.color.rgb = WHITE
    p.alignment  = PP_ALIGN.LEFT

    # Vertical centering tweak
    tf.margin_left   = Inches(0.45)
    tf.margin_top    = Inches(0.2)
    return bar_h


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(16), bold=False, color=DARK,
                font=BODY_FONT, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    r   = p.add_run()
    r.text = text
    r.font.name  = font
    r.font.size  = font_size
    r.font.bold  = bold
    r.font.italic= italic
    r.font.color.rgb = color
    p.alignment  = align
    return txb


def bullet_box(slide, items, left, top, width, height,
               font_size=Pt(15), color=DARK, indent_px=14,
               heading=None, heading_color=None):
    """
    items: list of strings or (str, level) tuples.
    level 0 = bullet, level 1 = sub-bullet.
    """
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True

    if heading:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = heading
        r.font.name  = BODY_FONT
        r.font.size  = Pt(17)
        r.font.bold  = True
        r.font.color.rgb = heading_color or MID_BL
        p.space_after = Pt(4)
        # next items on new paragraphs
        first = True
    else:
        first = True

    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if first and not heading:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        bullet = "  •  " if level == 0 else "      ‒  "
        r = p.add_run()
        r.text = bullet + text
        r.font.name  = BODY_FONT
        r.font.size  = font_size if level == 0 else Pt(font_size.pt - 1)
        r.font.color.rgb = color
        p.space_after = Pt(3)

    return txb


def accent_line(slide, top, color=MID_BL, thickness=Pt(2)):
    """Thin horizontal rule below title bar."""
    ln = slide.shapes.add_shape(1, 0, top, W, thickness)
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.width = 0


def two_col(slide, bar_h, left_items, right_items,
            left_head=None, right_head=None,
            font_size=Pt(15)):
    """Two-column bullet layout below the title bar."""
    margin  = Inches(0.45)
    col_w   = (W - margin * 3) / 2
    top     = bar_h + Inches(0.25)
    avail_h = H - top - Inches(0.2)

    bullet_box(slide, left_items,  margin,          top, col_w, avail_h,
               font_size=font_size, heading=left_head)
    bullet_box(slide, right_items, margin + col_w + margin, top, col_w, avail_h,
               font_size=font_size, heading=right_head)


def divider_slide(prs, section_title, subtitle=""):
    """Full-navy section divider."""
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.width = 0

    # Rule
    r = slide.shapes.add_shape(1, Inches(0.6), Inches(3.2), Inches(4), Pt(3))
    r.fill.solid(); r.fill.fore_color.rgb = MID_BL; r.line.width = 0

    add_textbox(slide, section_title,
                Inches(0.6), Inches(2.5), Inches(12), Inches(1.5),
                font_size=Pt(44), bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.6), Inches(3.7), Inches(11), Inches(1),
                    font_size=Pt(22), color=RGBColor(0xB0, 0xC8, 0xE8))
    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """01 — Title slide"""
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.width = 0

    # Accent stripe
    stripe = slide.shapes.add_shape(1, 0, Inches(5.8), W, Inches(0.08))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = MID_BL; stripe.line.width = 0

    add_textbox(slide,
                "Intelligent Grasp Failure Detection\nand Autonomous Recovery",
                Inches(0.7), Inches(1.4), Inches(11.5), Inches(2.5),
                font_size=Pt(40), bold=True, color=WHITE)

    add_textbox(slide,
                "UR10 Robot  ·  OAK-D Lite  ·  OnRobot RG2  ·  CLIP Vision",
                Inches(0.7), Inches(3.9), Inches(11), Inches(0.6),
                font_size=Pt(20), color=RGBColor(0xB0, 0xC8, 0xE8))

    add_textbox(slide,
                "Myu Wai Shin  ·  PDE3802 Final Year Project  ·  Middlesex University  ·  2025/26",
                Inches(0.7), Inches(6.1), Inches(11), Inches(0.7),
                font_size=Pt(15), color=RGBColor(0x90, 0xA8, 0xC8))
    return slide


def slide_introduction(prs):
    """02 — Introduction / problem statement"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Introduction")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "Pick-and-place is a fundamental robot task — but standard systems assume the pick always works.",
            "In practice, grasps fail silently: the robot carries an empty gripper to the place position.",
            "This work focuses on the under-studied problem of adding intelligence to detect and recover from failures.",
            "Key question: does the robot know it failed? What does it do about it?",
        ],
        right_items=[
            "Two-part pipeline:",
            ("Grasp failure detection — did the robot actually pick the object?", 1),
            ("Failure recovery — retry intelligently if it failed", 1),
            "Combined into a fully autonomous end-to-end pick-and-place loop.",
            "System flowchart shows the complete decision path: detect target → pick → verify → recover if failed.",
        ],
        left_head="Problem", right_head="This Project"
    )


def slide_aims(prs):
    """03 — Aims & Objectives"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Aims & Objectives")
    accent_line(slide, bar_h)

    margin = Inches(0.45)
    top = bar_h + Inches(0.3)
    col_w = W - margin * 2

    add_textbox(slide, "Aim: Develop an autonomous pick-and-place system that can detect grasp failures and recover.",
                margin, top, col_w, Inches(0.55),
                font_size=Pt(16), bold=True, color=NAVY)

    bullet_box(slide, [
        "YOLO object detection — locate objects and classify them (cube / cylinder)",
        "Pose estimation — transform detected object position to robot coordinate frame",
        "Read and interpret gripper IO signals — width (AI2 voltage) and contact force (DI8)",
        "Train a binary visual classifier (CLIP + SVM probe) to confirm object in gripper",
        "Implement 3-layer grasp failure detection: width → contact → visual (CLIP)",
        "Develop slip detection during transit using continuous gripper monitoring",
        "Integrate detection + recovery into a full end-to-end autonomous pipeline",
        "Evaluate against quantitative targets: TPR ≥ 85%, TNR ≥ 85%, FPR < 20%, task completion ≥ 70%",
    ],
    margin, top + Inches(0.65), col_w, Inches(4.8), font_size=Pt(15))


def slide_hardware(prs):
    """04 — Hardware"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Hardware")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "Universal Robots UR10",
            ("Industrial 6-DOF collaborative arm with built-in IO registers", 1),
            ("TCP/IP interface (ports 29999 and 30002) for robot control and sensor streaming", 1),
            ("Does not expose grip force directly — must be read via secondary interface", 1),
            "OnRobot RG2 Parallel Jaw Gripper",
            ("Width sensor via AI2 (analog voltage, calibrated 10.5–91 mm)", 1),
            ("Contact force limit via DI8 (digital signal HIGH when force limit reached)", 1),
            ("Gripper programs loaded as URP files via teach pendant", 1),
        ],
        right_items=[
            "OAK-D Lite (Luxonis)",
            ("RGB-D stereo camera — provides colour image and depth per pixel", 1),
            ("Eye-in-hand mount — wrist-mounted, moves with robot TCP", 1),
            ("Used for: object detection, pose estimation, and CLIP visual verification", 1),
            "Why not just one sensor?",
            ("Width alone can't distinguish a missed pick from a loose grip", 1),
            ("DI8 alone fires on partial contact", 1),
            ("Vision (CLIP) adds a redundant 3rd check — catches edge cases both sensors miss", 1),
        ],
        left_head="Robot + Gripper", right_head="Camera & Sensor Fusion"
    )


def slide_software(prs):
    """05 — Software & Libraries"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Software & Libraries")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "Python 3.10 — primary development language",
            "DepthAI SDK — OAK-D Lite camera driver and pipeline API",
            "OpenCV (cv2) — image processing, ArUco detection, display",
            "PyTorch + CUDA — GPU inference for CLIP and YOLO",
            "OpenAI CLIP (ViT-B/32) — visual verification classifier",
            "scikit-learn — SVM probe training on CLIP embeddings",
            "Ultralytics YOLOv8 / YOLOv5 — object detection and OBB",
        ],
        right_items=[
            "socket + struct — binary packet decoding from UR secondary interface",
            "threading — background sensor monitoring, parallel robot/camera loops",
            "numpy — coordinate transforms, linear algebra",
            "Pillow (PIL) — image preprocessing for CLIP",
            "autodistill + GroundingDINO — auto-annotation of training dataset",
            "OpenCV calibrateHandEye() — eye-in-hand calibration solver",
            "URScript via port 30002 — direct robot motion commands",
        ],
        left_head="Core Libraries", right_head="Communication & Tools"
    )


def slide_yolo_detection(prs):
    """06 — YOLO Object Detection"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "YOLO Object Detection")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "YOLOv5n, YOLOv8n, YOLOv8-OBB models trained on custom dataset",
            "V3 dataset: 2,525 images, cubes and cylinders, auto-annotated with GroundingDINO",
            "Deployed on OAK-D Lite using DepthAI pipeline (on-device inference)",
            "Depth estimate per detection from stereo disparity map",
        ],
        right_items=[
            "Modifications made for reliability:",
            ("Increased confidence threshold — filter out low-confidence detections", 1),
            ("NMS IoU threshold reduced — suppress overlapping duplicate boxes", 1),
            ("Temporal smoothing — only report detection if seen for 3 consecutive frames", 1),
            "OBB (Oriented Bounding Box) — gives object orientation for best grip angle",
            ("Standard box can't tell orientation of a cylinder or cube", 1),
            ("OBB wraps tightly and provides rotation angle", 1),
            ("Allows gripper to align to best pick orientation", 1),
        ],
        left_head="Training & Deployment", right_head="Key Modifications"
    )


def slide_yolo_limitations(prs):
    """07 — YOLO Limitations"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "YOLO — Limitations & Notes")
    accent_line(slide, bar_h)

    bullet_box(slide, [
        "OBB model trained on slow camera rotation sweep — good for orientation but limited viewpoints",
        "Mosaic augmentation during training caused confusion with backgrounds",
        "Low confidence on edge cases; temporal smoothing helps but adds latency",
        "No collision avoidance — detection does not account for robot arm obstructing view",
        "Works best on simple objects (cubes, cylinders); complex shapes need larger vision models",
        "Auto-annotation (GroundingDINO) can over-fire — ~10 false boxes/image at low threshold",
        "Not 100% accurate — failure detection pipeline adds a safety net on top of detection errors",
    ],
    Inches(0.45), bar_h + Inches(0.3), W - Inches(0.9), Inches(4.5), font_size=Pt(15))


def slide_pose_estimation(prs):
    """08 — Pose Estimation"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Pose Estimation")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "Goal: transform detected object pixel location into robot base-frame XYZ coordinates",
            "Camera configuration: Eye-in-hand — OAK-D Lite mounted on robot wrist (TCP)",
            "Calibration is critical — without it, camera and robot frames are misaligned",
            "Hand-eye calibration solves: camera → TCP rigid transform (R, t)",
            "Used OpenCV calibrateHandEye() with 5 solver methods (Tsai, Park, Horaud…)",
            "20–40 pose samples collected at varied positions and wrist angles",
            "Reprojection error < 2.0 px (excellent), < 5.0 px (acceptable)",
        ],
        right_items=[
            "ArUco marker placed on pick target for reliable 6-DoF pose detection",
            "PyStep: detect corners → solvePnP → tag pose in camera frame",
            "Transform chain:",
            ("Tag (camera frame) × R_cam2tcp → Tag (TCP frame)", 1),
            ("Tag (TCP frame) × FK (live robot pose) → Tag (base frame)", 1),
            ("Output: XYZ pick target in mm (robot coordinates)", 1),
            "Visual alignment HUD: live TCP crosshair vs. tag centre overlay for verification",
            "Limitation: calibration must be re-run if camera mount shifts",
        ],
        left_head="Calibration", right_head="Detection → Robot Frame"
    )


def slide_io_signals(prs):
    """09 — Reading IO Signals"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Reading Gripper IO Signals")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "All signals read from UR secondary interface — port 30002 binary stream",
            "Python struct module decodes binary packets in real time",
            "Background daemon thread updates sensor state continuously",
            "AI2 (Analog Input 2): gripper width voltage",
            ("Raw voltage 0–3.7 V → raw mm 0–110 mm", 1),
            ("Linear calibration applied (slope 1.405, offset -1.44)", 1),
            ("Fully open ≈ 91 mm   |   Fully closed ≈ 10.5 mm", 1),
        ],
        right_items=[
            "DI8 (Digital Input 8 / TDI1): contact force flag",
            ("Goes HIGH when gripper reaches force limit (object contacted)", 1),
            ("Used as a boolean: HIGH = contact confirmed, LOW = no contact", 1),
            "UR10 does not expose raw grip force — DI8 is the available signal",
            "Gripper state logic:",
            ("width ≥ 75 mm  →  gripper still open (close command failed)", 1),
            ("width < 11 mm  →  fully closed, missed object", 1),
            ("12 < width < 75 mm AND DI8 HIGH  →  object confirmed", 1),
        ],
        left_head="Communication", right_head="Signal Interpretation"
    )


def slide_failure_architecture(prs):
    """10 — Failure Detection Architecture"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Grasp Failure Detection — Architecture")
    accent_line(slide, bar_h)

    margin = Inches(0.45)
    top    = bar_h + Inches(0.3)

    add_textbox(slide, "Three independent detection layers — each is a hard gate:",
                margin, top, W - margin * 2, Inches(0.45),
                font_size=Pt(16), bold=True, color=NAVY)

    layers = [
        ("Layer 1", "Width Sensor (AI2)",     "Width in target range? Did the close command run?",   ACCENT),
        ("Layer 2", "Contact Signal (DI8)",   "DI8 HIGH? Force limit reached = object being gripped.", MID_BL),
        ("Layer 3", "CLIP Visual Check",      "Camera image classified as 'Holding' ≥ 75% confidence?", AMBER),
    ]

    box_w  = (W - margin * 4) / 3
    box_h  = Inches(2.2)
    box_top= top + Inches(0.55)

    for i, (tag, title, desc, col) in enumerate(layers):
        bx = margin + i * (box_w + margin)

        hdr = slide.shapes.add_shape(1, bx, box_top, box_w, Inches(0.5))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = col; hdr.line.width = 0

        tf = hdr.text_frame
        p  = tf.paragraphs[0]
        r  = p.add_run()
        r.text = f"{tag}  ·  {title}"
        r.font.name=BODY_FONT; r.font.size=Pt(13); r.font.bold=True
        r.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.margin_top = Inches(0.09)

        body = slide.shapes.add_shape(1, bx, box_top + Inches(0.5), box_w, box_h - Inches(0.5))
        body.fill.solid(); body.fill.fore_color.rgb = WHITE
        body.line.color.rgb = col; body.line.width = Pt(1.2)

        add_textbox(slide, desc,
                    bx + Inches(0.1), box_top + Inches(0.58),
                    box_w - Inches(0.2), box_h - Inches(0.65),
                    font_size=Pt(13), color=DARK)

    # Flow text
    flow_top = box_top + box_h + Inches(0.2)
    add_textbox(slide,
                "DESCEND  →  CLOSE  →  [Layer 1+2: width + DI8]  →  LIFT  →  [Post-lift width]  →  [Layer 3: CLIP]  →  TRANSFER",
                margin, flow_top, W - margin*2, Inches(0.5),
                font_size=Pt(13), color=MID_BL, bold=True)

    bullet_box(slide, [
        "Any layer failing triggers retry (up to MAX_RETRIES = 3)",
        "Each failure path logs a specific failure type to the terminal",
    ],
    margin, flow_top + Inches(0.55), W - margin*2, Inches(0.8), font_size=Pt(13))


def slide_slip_detection(prs):
    """11 — Slip Detection"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Slip Detection During Transit")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "Problem: object can slip OUT of gripper after grasp is confirmed — during movement",
            "Standard detection only checks at pick point — misses mid-transfer drops",
            "Solution: activate slip monitoring loop immediately after successful CLIP verification",
            "Loop-close URP: gripper re-closes every cycle during transit, compensating micro-slips",
            "Background thread monitors AI2 width continuously while robot moves",
        ],
        right_items=[
            "Slip event logic:",
            ("Had object (width > 11 mm after pick) = True", 1),
            ("Width drops below 11 mm during transit → slip detected", 1),
            ("_slip_detected flag set → pipeline pauses", 1),
            "On slip detected: CLIP re-verify triggered automatically",
            ("If CLIP still says 'Holding' → false alarm, continue", 1),
            ("If CLIP says 'Empty' → object lost, return and retry", 1),
            "This adds a 4th check not possible with gripper sensors alone",
        ],
        left_head="Problem & Approach", right_head="Slip Logic"
    )


def slide_clip(prs):
    """12 — CLIP Visual Verification"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "CLIP Visual Verification")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "CLIP (Contrastive Language–Image Pre-training) — OpenAI, ViT-B/32",
            "Encodes images and text into the same embedding space",
            "Similarity between image embedding and 'a gripper holding an object' vs 'an empty gripper'",
            "Originally zero-shot: no labelled gripper images needed",
            "Upgraded: trained a lightweight SVM probe on top of CLIP embeddings",
            ("Custom dataset: 600+ images of gripper holding / empty states", 1),
            ("Collected with OAK-D Lite camera at pick verification height", 1),
            ("SVM probe trained on CLIP features → 'Holding' / 'Empty'", 1),
        ],
        right_items=[
            "Integration in pipeline:",
            ("Robot lifts to verification height (200 mm above pick Z)", 1),
            ("Live camera frame captured — gripper region cropped", 1),
            ("CLIP encodes crop → SVM probe classifies", 1),
            ("Result displayed on monitor: label + confidence %", 1),
            ("Result auto-expires after 3 s — window resets to standby", 1),
            "Two stages:",
            ("'Grasp Verification' — after initial lift", 1),
            ("'Slip Re-verify' — after slip event detected during transit", 1),
            "Threshold: ≥ 75% confidence to pass; below = fail / retry",
        ],
        left_head="What CLIP Is", right_head="How It Is Used"
    )


def slide_recovery(prs):
    """13 — Recovery"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Recovery")
    accent_line(slide, bar_h)

    margin = Inches(0.45)
    top    = bar_h + Inches(0.3)

    add_textbox(slide, "Recovery strategy: re-approach and retry (up to MAX_RETRIES = 3)",
                margin, top, W - margin*2, Inches(0.45),
                font_size=Pt(16), bold=True, color=NAVY)

    bullet_box(slide, [
        "Layer 1 fail (width / DI8): open gripper → return to pick approach → close and recheck",
        "Layer 2 fail (post-lift width): open gripper → return to pick approach → retry",
        "Layer 3 fail (CLIP): open gripper → return to approach → retry from top",
        "Slip mid-transit: CLIP re-verify first → if still lost, return to pick and retry",
        "After MAX_RETRIES exhausted: pipeline reports failure and halts",
    ],
    margin, top + Inches(0.55), W - margin*2, Inches(2.5), font_size=Pt(15))

    add_textbox(slide, "⚠  Current Status — Work In Progress",
                margin, top + Inches(3.2), W - margin*2, Inches(0.45),
                font_size=Pt(15), bold=True, color=AMBER)

    bullet_box(slide, [
        "Detection of each failure type is fully implemented",
        "Re-approach and retry loop is implemented",
        "Advanced recovery strategies (e.g. adjusting pick pose, re-scanning) are planned but not yet integrated",
        "Limitation: calibration failure and collision avoidance are outside current scope",
    ],
    margin, top + Inches(3.75), W - margin*2, Inches(2.0), font_size=Pt(14))


def slide_pipeline(prs):
    """14 — End-to-End Pipeline"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "End-to-End Autonomous Pick-and-Place")
    accent_line(slide, bar_h)

    margin = Inches(0.45)
    top    = bar_h + Inches(0.25)

    steps = [
        ("1", "Detect",     "YOLO + DepthAI detect object. Depth gives approx Z."),
        ("2", "Locate",     "ArUco marker → pose estimation → robot XYZ coords."),
        ("3", "Approach",   "Robot moves to pick approach height (100 mm above)."),
        ("4", "Grip",       "Descend → close RG2 → settle 1.5 s."),
        ("5", "Check 1",    "Width + DI8: object in gripper? If fail → retry."),
        ("6", "Lift",       "Raise to CLIP verification height (200 mm above pick Z)."),
        ("7", "Check 2",    "Post-lift width: still holding? If not → slip during lift."),
        ("8", "CLIP",       "Visual verification: 'Holding' ≥ 75%? If fail → retry."),
        ("9", "Transit",    "Transfer to place position with slip monitor active."),
        ("10","Place",      "Lower → open gripper → retract → task complete."),
    ]

    box_w  = (W - margin * 2) / 5
    box_h  = Inches(1.3)
    row2_top = top + box_h + Inches(0.2)

    for i, (num, title, desc) in enumerate(steps):
        row   = i // 5
        col   = i % 5
        bx    = margin + col * box_w
        by    = top + row * (box_h + Inches(0.2))
        col_c = MID_BL if row == 0 else NAVY

        hdr = slide.shapes.add_shape(1, bx, by, box_w - Inches(0.05), Inches(0.38))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = col_c; hdr.line.width = 0
        tf = hdr.text_frame
        p  = tf.paragraphs[0]
        r  = p.add_run()
        r.text = f"{num}. {title}"
        r.font.name=BODY_FONT; r.font.size=Pt(13); r.font.bold=True
        r.font.color.rgb=WHITE
        p.alignment=PP_ALIGN.CENTER
        tf.margin_top=Inches(0.04)

        add_textbox(slide, desc, bx + Inches(0.05), by + Inches(0.4),
                    box_w - Inches(0.12), box_h - Inches(0.45),
                    font_size=Pt(11), color=DARK)

    note_top = top + 2*(box_h + Inches(0.2)) + Inches(0.1)
    add_textbox(slide, "✓  All stages run autonomously in a loop — no human intervention required between picks.",
                margin, note_top, W - margin*2, Inches(0.45),
                font_size=Pt(14), bold=True, color=ACCENT)


def slide_testing(prs):
    """15 — Testing & Evaluation"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Testing & Evaluation Targets")
    accent_line(slide, bar_h)

    margin = Inches(0.45)
    top    = bar_h + Inches(0.3)

    add_textbox(slide, "90-trial evaluation protocol — 30 successful grasps, 30 failed grasps, 30 slip events",
                margin, top, W - margin*2, Inches(0.45),
                font_size=Pt(15), bold=True, color=NAVY)

    rows = [
        ("Metric", "Formula", "Target", "Meaning"),
        ("True Positive Rate (TPR)", "TP / (TP + FN)", "≥ 85%", "Successful grasps correctly confirmed"),
        ("True Negative Rate (TNR)", "TN / (TN + FP)", "≥ 85%", "Failed grasps correctly identified"),
        ("False Positive Rate (FPR)", "FP / (FP + TN)", "< 20%", "System says success when grasp failed"),
        ("Overall Accuracy",         "(TP + TN) / 90", "≥ 85%", "Combined correct verdicts across all trials"),
        ("Recovery Success Rate",    "Recoveries succeeded / attempted", "≥ 70%", "Failed grasps resolved by recovery"),
        ("End-to-end Completion",    "Successful placements / 90", "≥ 70%", "Overall system performance"),
    ]

    col_widths = [Inches(2.8), Inches(2.5), Inches(1.2), Inches(5.8)]
    row_h      = Inches(0.52)
    tbl_top    = top + Inches(0.6)

    for r_i, row in enumerate(rows):
        bg_c = NAVY if r_i == 0 else (LIGHT if r_i % 2 == 0 else WHITE)
        txt_c= WHITE if r_i == 0 else DARK

        x = margin
        for c_i, (cell, cw) in enumerate(zip(row, col_widths)):
            rect = slide.shapes.add_shape(1, x, tbl_top + r_i*row_h, cw - Inches(0.02), row_h)
            rect.fill.solid(); rect.fill.fore_color.rgb = bg_c
            rect.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); rect.line.width = Pt(0.5)

            add_textbox(slide, cell, x + Inches(0.07), tbl_top + r_i*row_h + Inches(0.06),
                        cw - Inches(0.15), row_h - Inches(0.05),
                        font_size=Pt(12 if r_i > 0 else 12),
                        bold=(r_i == 0), color=txt_c)
            x += cw


def slide_results_placeholder(prs):
    """16 — Results (placeholder)"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Results")
    accent_line(slide, bar_h)

    margin = Inches(0.45)
    top    = bar_h + Inches(0.35)

    add_textbox(slide, "[INSERT RESULTS TABLE / CHARTS HERE]",
                margin, top, W - margin*2, Inches(0.55),
                font_size=Pt(18), bold=True, color=AMBER)

    bullet_box(slide, [
        "Report measured TPR, TNR, FPR, Overall Accuracy, Recovery Rate, Task Completion Rate",
        "Show comparison: detection-only pipeline vs. full detection + recovery pipeline",
        "Include demo video screenshots or frames showing: successful pick, detected failure, recovery",
        "Highlight any cases where CLIP disagreed with width/DI8 — discuss sensor fusion benefit",
    ],
    margin, top + Inches(0.65), W - margin*2, Inches(3.5), font_size=Pt(15))

    add_textbox(slide, "→ Add your actual measured values and video evidence here before presenting.",
                margin, top + Inches(4.3), W - margin*2, Inches(0.45),
                font_size=Pt(14), italic=True, color=MID_BL)


def slide_limitations(prs):
    """17 — Limitations"""
    slide = blank_slide(prs)
    bg_rect(slide)
    bar_h = title_bar(slide, "Limitations & Future Work")
    accent_line(slide, bar_h)

    two_col(slide, bar_h,
        left_items=[
            "No collision avoidance — robot does not model obstacles in workspace",
            "Calibration drift — hand-eye calibration must be re-run if camera shifts",
            "CLIP accuracy not 100% — uncertain predictions on poor lighting or extreme crops",
            "Gripper force not readable from Python — only binary DI8 contact flag available",
            "Recovery re-approach assumes object is still in the same position",
            "No object re-detection if object moves between attempts",
        ],
        right_items=[
            "Advanced recovery strategies (adjust pick pose, rotate approach angle)",
            "Replace URP-based gripper control with direct force torque sensing",
            "Extend CLIP training dataset for more challenging lighting/backgrounds",
            "Add object re-detection before each retry attempt",
            "Generalise to arbitrary objects (beyond cubes and cylinders)",
            "Integrate with motion planning for collision-aware trajectories",
        ],
        left_head="Current Limitations", right_head="Future Work"
    )


def slide_conclusion(prs):
    """18 — Conclusion"""
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.width = 0

    r = slide.shapes.add_shape(1, Inches(0.6), Inches(2.8), Inches(6), Pt(3))
    r.fill.solid(); r.fill.fore_color.rgb = MID_BL; r.line.width = 0

    add_textbox(slide, "Summary",
                Inches(0.6), Inches(1.1), Inches(12), Inches(0.8),
                font_size=Pt(40), bold=True, color=WHITE)

    bullet_box(slide, [
        "Built a 3-layer grasp failure detection pipeline: width + DI8 + CLIP visual",
        "Implemented slip detection via continuous gripper monitoring during transit",
        "Integrated CLIP (ViT-B/32) + custom SVM probe for vision-based grasp confirmation",
        "Developed end-to-end autonomous pick-and-place with detection and recovery loop",
        "Evaluation targets: TPR ≥ 85%, FPR < 20%, Task completion ≥ 70%",
    ],
    Inches(0.6), Inches(3.0), Inches(12), Inches(3.5),
    font_size=Pt(17), color=RGBColor(0xD8, 0xE8, 0xF5))

    add_textbox(slide,
                "Myu Wai Shin  ·  PDE3802 Final Year Project  ·  Middlesex University  ·  2025/26",
                Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
                font_size=Pt(13), color=RGBColor(0x80, 0x9A, 0xB8))


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    # Section dividers + content slides
    slide_title(prs)                                       # 01
    slide_introduction(prs)                                # 02
    slide_aims(prs)                                        # 03
    divider_slide(prs, "Hardware & Software",
                  "What the system is built on")           # 04
    slide_hardware(prs)                                    # 05
    slide_software(prs)                                    # 06
    divider_slide(prs, "Object Detection",
                  "YOLO + OBB on OAK-D Lite")              # 07
    slide_yolo_detection(prs)                              # 08
    slide_yolo_limitations(prs)                            # 09
    divider_slide(prs, "Pose Estimation",
                  "Eye-in-hand calibration → robot frame") # 10
    slide_pose_estimation(prs)                             # 11
    divider_slide(prs, "Failure Detection Pipeline",
                  "3-layer grasp verification + slip")     # 12
    slide_io_signals(prs)                                  # 13
    slide_failure_architecture(prs)                        # 14
    slide_slip_detection(prs)                              # 15
    slide_clip(prs)                                        # 16
    slide_recovery(prs)                                    # 17
    slide_pipeline(prs)                                    # 18
    divider_slide(prs, "Testing & Results",
                  "Evaluation protocol and outcomes")      # 19
    slide_testing(prs)                                     # 20
    slide_results_placeholder(prs)                         # 21
    slide_limitations(prs)                                 # 22
    slide_conclusion(prs)                                  # 23

    out = "FYP_Presentation.pptx"
    prs.save(out)
    print(f"\n  Saved: {out}  ({len(prs.slides)} slides)\n")


if __name__ == "__main__":
    main()
