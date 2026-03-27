"""
build_presentation_v2.py
------------------------
FYP Presentation — visual-first redesign.
Run:  python build_presentation_v2.py
Out:  FYP_Presentation_v2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
BLUE   = RGBColor(0x2D, 0x5F, 0x8A)
LBLUE  = RGBColor(0xD8, 0xE8, 0xF5)
LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1C, 0x1C, 0x1C)
GREEN  = RGBColor(0x1E, 0x7E, 0x45)
AMBER  = RGBColor(0xC8, 0x72, 0x00)
MED    = RGBColor(0x55, 0x55, 0x55)

W = Inches(13.33)
H = Inches(7.5)
TF = "Calibri"

def prs_new():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ── Primitive helpers ─────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=LIGHT, line_color=None, line_w=0):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.width = 0
    return s

def txt(slide, text, l, t, w, h, size=Pt(16), bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = TF; r.font.size = size; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    p.alignment = align
    return b

def multiline_txt(slide, lines, l, t, w, h, size=Pt(15), color=DARK,
                  bold=False, line_spacing=Pt(4)):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.name = TF; r.font.size = size; r.font.bold = bold
        r.font.color.rgb = color
        p.space_after = line_spacing
    return b

def bg(slide, color=LIGHT):
    rect(slide, 0, 0, W, H, fill=color)

def topbar(slide, title_text, bar_h=Inches(1.1)):
    rect(slide, 0, 0, W, bar_h, fill=NAVY)
    txt(slide, title_text,
        Inches(0.5), Inches(0.15), W - Inches(1), bar_h - Inches(0.1),
        size=Pt(30), bold=True, color=WHITE)
    rect(slide, 0, bar_h, W, Pt(3), fill=BLUE)
    return bar_h

def tagbox(slide, label, l, t, w, h, bg_color=BLUE, text_color=WHITE, size=Pt(13)):
    rect(slide, l, t, w, h, fill=bg_color)
    txt(slide, label, l + Inches(0.08), t + Inches(0.07),
        w - Inches(0.16), h - Inches(0.1),
        size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def stat_box(slide, number, label, l, t, w=Inches(2.4), h=Inches(1.6),
             bg_color=NAVY):
    rect(slide, l, t, w, h, fill=bg_color)
    txt(slide, number, l, t + Inches(0.18), w, Inches(0.8),
        size=Pt(34), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, label, l, t + Inches(0.92), w, Inches(0.55),
        size=Pt(12), color=LBLUE, align=PP_ALIGN.CENTER)

def divider(prs, heading, sub=""):
    slide = blank(prs)
    rect(slide, 0, 0, W, H, fill=NAVY)
    rect(slide, Inches(0.6), Inches(3.05), Inches(5), Pt(4), fill=BLUE)
    txt(slide, heading, Inches(0.6), Inches(1.8), Inches(12), Inches(1.4),
        size=Pt(46), bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.6), Inches(3.35), Inches(11), Inches(0.8),
            size=Pt(20), color=RGBColor(0xA0,0xBE,0xDE))
    return slide

# ════════════════════════════════════════════════════════════════════════════
# INDIVIDUAL SLIDES
# ════════════════════════════════════════════════════════════════════════════

def s01_title(prs):
    slide = blank(prs)
    rect(slide, 0, 0, W, H, fill=NAVY)
    rect(slide, 0, Inches(5.6), W, Inches(0.07), fill=BLUE)

    txt(slide, "Intelligent Grasp Failure\nDetection and Autonomous Recovery",
        Inches(0.7), Inches(0.9), Inches(11.5), Inches(3.2),
        size=Pt(42), bold=True, color=WHITE)

    txt(slide, "UR10   ·   OnRobot RG2   ·   OAK-D Lite   ·   CLIP Vision",
        Inches(0.7), Inches(4.0), Inches(11), Inches(0.6),
        size=Pt(19), color=RGBColor(0xA0,0xBE,0xDE))

    txt(slide, "Myu Wai Shin  ·  PDE3802 Final Year Project  ·  Middlesex University  ·  2025/26",
        Inches(0.7), Inches(6.6), Inches(11), Inches(0.6),
        size=Pt(14), color=RGBColor(0x80,0xA0,0xC0))


def s02_intro(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Introduction")

    # Left — problem statement big text
    rect(slide, Inches(0.45), bar_h + Inches(0.25),
         Inches(6.2), H - bar_h - Inches(0.5), fill=WHITE,
         line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))

    txt(slide, "The Problem",
        Inches(0.65), bar_h + Inches(0.45), Inches(5.8), Inches(0.5),
        size=Pt(18), bold=True, color=BLUE)
    txt(slide,
        "Robotic pick-and-place systems assume the grasp works.\n\n"
        "When it fails, the robot carries an empty gripper — unaware.\n\n"
        "This is a critical gap: the system does not know it failed.",
        Inches(0.65), bar_h + Inches(0.95), Inches(5.8), Inches(3.5),
        size=Pt(17), color=DARK)

    # Right — this project
    rect(slide, Inches(7.1), bar_h + Inches(0.25),
         Inches(5.75), H - bar_h - Inches(0.5), fill=NAVY)

    txt(slide, "This Project",
        Inches(7.3), bar_h + Inches(0.45), Inches(5.3), Inches(0.5),
        size=Pt(18), bold=True, color=LBLUE)

    steps = [
        "① Detect whether the grasp succeeded",
        "② If it failed — detect why and recover",
        "③ Loop: fully autonomous pick-and-place",
    ]
    top = bar_h + Inches(0.95)
    for s in steps:
        txt(slide, s, Inches(7.3), top, Inches(5.3), Inches(0.8),
            size=Pt(16), color=WHITE)
        top += Inches(0.95)

    txt(slide, "Focus: adding intelligence to what happens after a failure — "
               "not just training the robot to pick.",
        Inches(7.3), top + Inches(0.1), Inches(5.3), Inches(1.1),
        size=Pt(14), color=RGBColor(0xA0,0xBE,0xDE), italic=True)


def s03_aims(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Aims & Objectives")

    aim_top = bar_h + Inches(0.25)
    rect(slide, Inches(0.45), aim_top, W - Inches(0.9), Inches(0.62), fill=BLUE)
    txt(slide, "Aim: Develop an autonomous pick-and-place system capable of detecting grasp failures and recovering.",
        Inches(0.62), aim_top + Inches(0.1), W - Inches(1.3), Inches(0.45),
        size=Pt(15), bold=True, color=WHITE)

    objs = [
        ("1", "YOLO Object Detection",        "Locate and classify objects (cube / cylinder) in the workspace"),
        ("2", "Pose Estimation",              "Transform pixel detection → robot XYZ coordinates"),
        ("3", "Gripper IO Signals",           "Read width (AI2) and contact force (DI8) from RG2 in real time"),
        ("4", "Visual Classifier (CLIP)",     "Linear probe on CLIP embeddings: confirm object in gripper"),
        ("5", "3-Layer Failure Detection",    "Width → Contact → CLIP: redundant independent checks"),
        ("6", "Slip Detection",               "Monitor grip during transit; detect mid-move object loss"),
        ("7", "Recovery Integration",         "Retry loop: return to pick, re-approach, re-verify"),
    ]

    col_w = (W - Inches(0.9)) / 2
    for i, (num, title, desc) in enumerate(objs):
        col = i % 2
        row = i // 2
        l   = Inches(0.45) + col * (col_w + Inches(0.05))
        t   = aim_top + Inches(0.75) + row * Inches(1.08)

        rect(slide, l, t, col_w - Inches(0.05), Inches(0.98),
             fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
        rect(slide, l, t, Inches(0.42), Inches(0.98), fill=NAVY)
        txt(slide, num, l, t + Inches(0.18), Inches(0.42), Inches(0.55),
            size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, title, l + Inches(0.5), t + Inches(0.06),
            col_w - Inches(0.62), Inches(0.38),
            size=Pt(13), bold=True, color=NAVY)
        txt(slide, desc, l + Inches(0.5), t + Inches(0.44),
            col_w - Inches(0.62), Inches(0.48),
            size=Pt(12), color=MED)


def s04_hardware(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Hardware")

    cards = [
        (NAVY,  "UR10\nRobot Arm",
         ["6-DOF industrial collaborative arm",
          "TCP/IP interface (ports 29999 & 30002)",
          "IO registers: AI2 (analog) + DI8 (digital)"]),
        (BLUE,  "OnRobot RG2\nGripper",
         ["Parallel jaw — max 110 mm width",
          "AI2: width voltage (calibrated ±1 mm)",
          "DI8: HIGH when force limit reached"]),
        (RGBColor(0x1E,0x6E,0x45), "OAK-D Lite\nCamera",
         ["RGB-D stereo — colour + depth per pixel",
          "Eye-in-hand: wrist-mounted, moves with TCP",
          "Detection · Pose · CLIP verification"]),
    ]

    card_w = (W - Inches(1.8)) / 3
    for i, (col, title, pts) in enumerate(cards):
        l = Inches(0.45) + i * (card_w + Inches(0.45))
        t = bar_h + Inches(0.3)
        rect(slide, l, t, card_w, Inches(1.0), fill=col)
        txt(slide, title, l + Inches(0.15), t + Inches(0.1),
            card_w - Inches(0.3), Inches(0.85),
            size=Pt(17), bold=True, color=WHITE)
        body_t = t + Inches(1.05)
        rect(slide, l, body_t, card_w, Inches(4.5),
             fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
        for j, pt in enumerate(pts):
            txt(slide, "• " + pt, l + Inches(0.15), body_t + Inches(0.2) + j * Inches(0.75),
                card_w - Inches(0.3), Inches(0.7),
                size=Pt(14), color=DARK)

    # Bottom: why fusion?
    bot_t = bar_h + Inches(5.65)
    rect(slide, Inches(0.45), bot_t, W - Inches(0.9), Inches(0.72), fill=LBLUE)
    txt(slide,
        "Why width + force + vision?  — Each sensor catches failures the others miss. "
        "Width alone can't confirm contact; DI8 fires on partial touch; CLIP adds a visual ground truth.",
        Inches(0.65), bot_t + Inches(0.1), W - Inches(1.3), Inches(0.55),
        size=Pt(14), color=NAVY)


def s05_software(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Software & Libraries")

    groups = [
        ("Robot Control",       NAVY,  ["socket + struct — binary UR packet decoding",
                                         "URScript via port 30002 — motion commands",
                                         "URP files (teach pendant) — gripper control"]),
        ("Vision & Detection",  BLUE,  ["DepthAI SDK — OAK-D Lite pipeline",
                                         "OpenCV — ArUco, image processing, display",
                                         "Ultralytics YOLOv8 / YOLOv5 OBB"]),
        ("AI & Classification", RGBColor(0x1E,0x6E,0x45),
                               ["PyTorch + CUDA — GPU inference",
                                "OpenAI CLIP ViT-B/32 — image embeddings",
                                "scikit-learn — Logistic Regression (linear probe)"]),
        ("Data & Tools",        RGBColor(0x5A,0x3A,0x8A),
                               ["numpy + Pillow — transforms, preprocessing",
                                "GroundingDINO + autodistill — auto-annotation",
                                "OpenCV calibrateHandEye() — eye-in-hand solver"]),
    ]

    gw = (W - Inches(1.8)) / 2
    gh = Inches(2.7)
    for i, (title, col, pts) in enumerate(groups):
        row = i // 2; c = i % 2
        l = Inches(0.45) + c * (gw + Inches(0.9))
        t = bar_h + Inches(0.3) + row * (gh + Inches(0.25))
        rect(slide, l, t, gw, Inches(0.45), fill=col)
        txt(slide, title, l + Inches(0.15), t + Inches(0.06),
            gw - Inches(0.3), Inches(0.35),
            size=Pt(14), bold=True, color=WHITE)
        rect(slide, l, t + Inches(0.45), gw, gh - Inches(0.45),
             fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
        for j, pt in enumerate(pts):
            txt(slide, "• " + pt,
                l + Inches(0.15), t + Inches(0.6) + j * Inches(0.62),
                gw - Inches(0.3), Inches(0.55),
                size=Pt(13), color=DARK)


def s06_yolo(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "YOLO Object Detection")

    # Left panel
    lw = Inches(5.8)
    rect(slide, Inches(0.45), bar_h + Inches(0.25), lw, H - bar_h - Inches(0.5),
         fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
    txt(slide, "Training & Deployment",
        Inches(0.65), bar_h + Inches(0.45), lw - Inches(0.4), Inches(0.42),
        size=Pt(16), bold=True, color=BLUE)
    multiline_txt(slide,
        ["Models: YOLOv5n, YOLOv8n, YOLOv8-OBB",
         "Dataset V3: 2,525 images — cubes & cylinders",
         "Auto-annotated with GroundingDINO",
         "Deployed on OAK-D Lite via DepthAI pipeline",
         "Depth estimate from stereo disparity per detection"],
        Inches(0.65), bar_h + Inches(0.95), lw - Inches(0.4), Inches(3.0),
        size=Pt(15))

    # Right: modifications as 3 highlight boxes
    rstart = Inches(0.45) + lw + Inches(0.4)
    rw = W - rstart - Inches(0.45)
    txt(slide, "Key Modifications",
        rstart, bar_h + Inches(0.45), rw, Inches(0.42),
        size=Pt(16), bold=True, color=BLUE)

    mods = [
        (NAVY,  "Confidence Threshold ↑",
                "Raised threshold to filter out weak detections"),
        (BLUE,  "NMS IoU ↓",
                "Suppresses duplicate overlapping boxes"),
        (RGBColor(0x1E,0x6E,0x45), "Temporal Smoothing",
                "Only reports object if seen for 3 consecutive frames"),
        (RGBColor(0x5A,0x3A,0x8A), "OBB — Orientation",
                "Oriented Bounding Box gives rotation angle → best grip"),
    ]
    mt = bar_h + Inches(0.95)
    for col, title, desc in mods:
        rect(slide, rstart, mt, rw, Inches(1.2), fill=col)
        txt(slide, title, rstart + Inches(0.15), mt + Inches(0.1),
            rw - Inches(0.3), Inches(0.45),
            size=Pt(14), bold=True, color=WHITE)
        txt(slide, desc, rstart + Inches(0.15), mt + Inches(0.55),
            rw - Inches(0.3), Inches(0.55),
            size=Pt(13), color=LBLUE)
        mt += Inches(1.28)


def s07_pose(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Pose Estimation")

    # Step boxes
    steps = [
        ("Camera\nCalibration",   "Intrinsics via\ncharuco board"),
        ("Hand-Eye\nCalibration", "20–40 poses → OpenCV\ncalibrateHandEye()\nsolves camera→TCP"),
        ("ArUco\nDetection",      "6×6 marker on target →\n6-DoF pose in\ncamera frame"),
        ("Coord\nTransform",      "Camera → TCP → Base\nframe using FK\n(live robot pose)"),
        ("Robot\nXYZ Output",     "Pick target in mm\nin robot base frame"),
    ]

    sw = (W - Inches(0.9)) / len(steps) - Inches(0.1)
    for i, (title, desc) in enumerate(steps):
        l = Inches(0.45) + i * (sw + Inches(0.1))
        t = bar_h + Inches(0.3)
        # header
        rect(slide, l, t, sw, Inches(0.9), fill=NAVY)
        txt(slide, title, l, t + Inches(0.06), sw, Inches(0.82),
            size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # body
        rect(slide, l, t + Inches(0.9), sw, Inches(2.4),
             fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
        txt(slide, desc, l + Inches(0.1), t + Inches(1.0),
            sw - Inches(0.2), Inches(2.2),
            size=Pt(13), color=DARK, align=PP_ALIGN.CENTER)
        # arrow (not last)
        if i < len(steps) - 1:
            txt(slide, "→", l + sw + Inches(0.0), t + Inches(0.35),
                Inches(0.12), Inches(0.5),
                size=Pt(20), bold=True, color=MED, align=PP_ALIGN.CENTER)

    # Key point
    kt = bar_h + Inches(3.5)
    rect(slide, Inches(0.45), kt, W - Inches(0.9), Inches(0.62), fill=LBLUE)
    txt(slide, "Eye-in-hand: camera moves with the wrist — calibration gives a fixed camera→TCP offset, "
               "valid for all robot poses.  Reprojection error achieved: < 2.0 px.",
        Inches(0.65), kt + Inches(0.1), W - Inches(1.3), Inches(0.45),
        size=Pt(14), color=NAVY)

    # Limitation row
    lt = kt + Inches(0.75)
    rect(slide, Inches(0.45), lt, W - Inches(0.9), Inches(2.3),
         fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
    txt(slide, "Limitations",
        Inches(0.65), lt + Inches(0.1), Inches(3), Inches(0.38),
        size=Pt(14), bold=True, color=BLUE)
    multiline_txt(slide,
        ["• Calibration must be re-run if camera mount shifts",
         "• Assumes ArUco marker is placed on or near the pick target",
         "• Accuracy depends on board detection quality and number of poses collected"],
        Inches(0.65), lt + Inches(0.5), W - Inches(1.3), Inches(1.7),
        size=Pt(14))


def s08_io(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Reading Gripper IO Signals")

    # Two signal cards
    signals = [
        (NAVY, "AI2 — Width (Analog)",
         "Gripper finger spread in mm",
         ["Voltage read from UR secondary interface (port 30002)",
          "Raw voltage → raw mm → calibrated mm",
          "Linear calibration:  slope 1.405 · offset -1.44",
          "Fully open ≈ 91 mm   |   Fully closed ≈ 10.5 mm",
          "Accuracy: ± 1 mm"]),
        (BLUE, "DI8 — Contact Force (Digital)",
         "ON/OFF flag when force limit reached",
         ["Digital register read from same binary stream",
          "HIGH = gripper reached programmed force limit",
          "Acts as a contact confirmation signal",
          "UR10 does not expose raw force value — DI8 is the available signal",
          "Background daemon thread monitors both signals at ~25 Hz"]),
    ]

    cw = (W - Inches(1.35)) / 2
    for i, (col, title, sub, pts) in enumerate(signals):
        l = Inches(0.45) + i * (cw + Inches(0.45))
        t = bar_h + Inches(0.3)
        rect(slide, l, t, cw, Inches(0.9), fill=col)
        txt(slide, title, l + Inches(0.15), t + Inches(0.06),
            cw - Inches(0.3), Inches(0.45),
            size=Pt(16), bold=True, color=WHITE)
        txt(slide, sub, l + Inches(0.15), t + Inches(0.52),
            cw - Inches(0.3), Inches(0.3),
            size=Pt(13), color=LBLUE)
        rect(slide, l, t + Inches(0.9), cw, Inches(4.5),
             fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
        for j, pt in enumerate(pts):
            txt(slide, "• " + pt,
                l + Inches(0.15), t + Inches(1.05) + j * Inches(0.72),
                cw - Inches(0.3), Inches(0.65),
                size=Pt(14), color=DARK)

    # Decision table
    dt = bar_h + Inches(5.55)
    decisions = [
        ("width ≥ 75 mm", "Gripper still open — close command failed"),
        ("width < 11 mm", "Fully closed — missed object"),
        ("12–75 mm & DI8 HIGH", "Object confirmed  ✓"),
    ]
    dw = (W - Inches(0.9)) / 3
    cols_d = [AMBER, AMBER, GREEN]
    for i, (cond, meaning) in enumerate(decisions):
        l = Inches(0.45) + i * dw
        rect(slide, l, dt, dw - Inches(0.05), Inches(0.82), fill=cols_d[i])
        txt(slide, cond, l + Inches(0.1), dt + Inches(0.05),
            dw - Inches(0.2), Inches(0.35),
            size=Pt(13), bold=True, color=WHITE)
        txt(slide, meaning, l + Inches(0.1), dt + Inches(0.4),
            dw - Inches(0.2), Inches(0.35),
            size=Pt(12), color=WHITE)


def s09_detection_arch(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Grasp Failure Detection — 3 Layers")

    layers = [
        (NAVY,  "Layer 1",  "Width + DI8",
         "Width in valid range?\n& DI8 HIGH (contact)?",
         "Checks immediately after gripping"),
        (BLUE,  "Layer 2",  "Post-lift Width",
         "Width still > 11 mm\nafter lifting?",
         "Catches objects slipped during lift"),
        (RGBColor(0x1E,0x6E,0x45), "Layer 3", "CLIP Visual",
         "Camera: 'Holding'\n≥ 75% confidence?",
         "Visual ground truth — independent of sensors"),
    ]

    lw = (W - Inches(1.8)) / 3
    lh = Inches(4.0)
    lt = bar_h + Inches(0.3)

    for i, (col, tag, name, question, note) in enumerate(layers):
        l = Inches(0.45) + i * (lw + Inches(0.45))
        rect(slide, l, lt, lw, Inches(0.55), fill=col)
        txt(slide, f"{tag}  ·  {name}", l + Inches(0.12), lt + Inches(0.08),
            lw - Inches(0.24), Inches(0.42),
            size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        rect(slide, l, lt + Inches(0.55), lw, Inches(2.0),
             fill=WHITE, line_color=col, line_w=Pt(2))
        txt(slide, "Checks:", l + Inches(0.15), lt + Inches(0.65),
            lw - Inches(0.3), Inches(0.35),
            size=Pt(12), bold=True, color=col)
        txt(slide, question, l + Inches(0.15), lt + Inches(1.02),
            lw - Inches(0.3), Inches(1.3),
            size=Pt(16), bold=True, color=DARK)

        rect(slide, l, lt + Inches(2.55), lw, Inches(0.72), fill=LBLUE)
        txt(slide, note, l + Inches(0.12), lt + Inches(2.62),
            lw - Inches(0.24), Inches(0.6),
            size=Pt(13), color=NAVY, italic=True)

        # Fail / pass labels
        rect(slide, l, lt + Inches(3.35), lw / 2 - Inches(0.05), Inches(0.5),
             fill=RGBColor(0xC8,0x35,0x35))
        txt(slide, "FAIL → retry", l + Inches(0.05), lt + Inches(3.38),
            lw / 2 - Inches(0.1), Inches(0.38),
            size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, l + lw / 2 + Inches(0.05), lt + Inches(3.35),
             lw / 2 - Inches(0.05), Inches(0.5), fill=GREEN)
        txt(slide, "PASS → next", l + lw / 2 + Inches(0.1), lt + Inches(3.38),
            lw / 2 - Inches(0.1), Inches(0.38),
            size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)

    # Flow bar at bottom
    flow_t = lt + Inches(3.95)
    txt(slide,
        "DESCEND  →  GRIP  →  [L1]  →  LIFT  →  [L2]  →  [L3: CLIP]  →  TRANSFER  →  [Slip?]  →  PLACE",
        Inches(0.45), flow_t, W - Inches(0.9), Inches(0.5),
        size=Pt(14), bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def s10_clip(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "CLIP Visual Verification")

    # What is CLIP panel
    lw = Inches(5.5)
    rect(slide, Inches(0.45), bar_h + Inches(0.25), lw, H - bar_h - Inches(0.5),
         fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
    txt(slide, "What is CLIP?",
        Inches(0.65), bar_h + Inches(0.45), lw, Inches(0.42),
        size=Pt(16), bold=True, color=BLUE)
    txt(slide,
        "CLIP (OpenAI, ViT-B/32) encodes images and text into "
        "the same vector space — making it possible to compare an image "
        "directly to a text description.\n\n"
        "Image crop of gripper  →  CLIP  →  512-dim feature vector\n\n"
        "A Logistic Regression linear probe is trained on top of these "
        "features using ~600 real gripper images:\n\n"
        "  • Holding  (object between fingers)\n"
        "  • Empty  (no object)\n\n"
        "Result: label + confidence % in < 200 ms.",
        Inches(0.65), bar_h + Inches(0.95), lw - Inches(0.4), Inches(4.5),
        size=Pt(15), color=DARK)

    # How it works in pipeline
    rstart = Inches(0.45) + lw + Inches(0.4)
    rw = W - rstart - Inches(0.45)
    txt(slide, "In the Pipeline",
        rstart, bar_h + Inches(0.45), rw, Inches(0.42),
        size=Pt(16), bold=True, color=BLUE)

    pipeline_steps = [
        (NAVY,  "①  Lift to CLIP height",  "200 mm above pick position"),
        (NAVY,  "②  Flush camera buffer",  "0.6 s settle — ensures fresh, sharp frame"),
        (NAVY,  "③  Crop gripper region",  "Fixed crop centred on gripper tip"),
        (NAVY,  "④  CLIP → Linear probe",  "'Holding' or 'Empty' + confidence"),
        (GREEN, "⑤  ≥ 75%  'Holding'",     "✓ Pass — continue to transfer"),
        (RGBColor(0xC8,0x35,0x35), "⑤  < 75% or 'Empty'", "✗ Fail — open, return, retry"),
    ]

    st = bar_h + Inches(0.95)
    for col, title, sub in pipeline_steps:
        rect(slide, rstart, st, rw, Inches(0.8), fill=col)
        txt(slide, title, rstart + Inches(0.12), st + Inches(0.06),
            rw - Inches(0.24), Inches(0.35),
            size=Pt(13), bold=True, color=WHITE)
        txt(slide, sub, rstart + Inches(0.12), st + Inches(0.42),
            rw - Inches(0.24), Inches(0.32),
            size=Pt(12), color=LBLUE)
        st += Inches(0.87)


def s11_slip(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Slip Detection During Transfer")

    # Problem
    rect(slide, Inches(0.45), bar_h + Inches(0.25),
         W - Inches(0.9), Inches(0.9), fill=NAVY)
    txt(slide,
        "Problem:  gripper sensors only checked at pick point — object can still slip during transit "
        "with no detection.",
        Inches(0.65), bar_h + Inches(0.35), W - Inches(1.3), Inches(0.72),
        size=Pt(15), bold=True, color=WHITE)

    # Two columns: mechanism + logic
    cw = (W - Inches(1.35)) / 2
    ct = bar_h + Inches(1.3)
    ch = H - ct - Inches(0.3)

    # Mechanism
    rect(slide, Inches(0.45), ct, cw, ch,
         fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
    txt(slide, "How It Works",
        Inches(0.65), ct + Inches(0.15), cw, Inches(0.4),
        size=Pt(16), bold=True, color=BLUE)
    multiline_txt(slide,
        ["Loop-close URP: gripper re-closes continuously during transit",
         "",
         "Background thread reads AI2 width at ~25 Hz",
         "",
         "If width drops below 11 mm mid-transit → slip flag set",
         "",
         "CLIP re-verify triggered immediately"],
        Inches(0.65), ct + Inches(0.62), cw - Inches(0.4), ch - Inches(0.8),
        size=Pt(15))

    # Logic
    rx = Inches(0.45) + cw + Inches(0.45)
    rect(slide, rx, ct, cw, ch,
         fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
    txt(slide, "Slip Response",
        rx + Inches(0.2), ct + Inches(0.15), cw, Inches(0.4),
        size=Pt(16), bold=True, color=BLUE)

    outcomes = [
        (BLUE,  "CLIP re-verify",  "Triggers after slip flag set"),
        (GREEN, "Still 'Holding'", "False alarm → continue to place"),
        (RGBColor(0xC8,0x35,0x35), "'Empty'", "Object lost → return to pick, retry"),
    ]
    ot = ct + Inches(0.65)
    for col, label, sub in outcomes:
        rect(slide, rx + Inches(0.2), ot, cw - Inches(0.4), Inches(1.1), fill=col)
        txt(slide, label, rx + Inches(0.35), ot + Inches(0.08),
            cw - Inches(0.55), Inches(0.45),
            size=Pt(15), bold=True, color=WHITE)
        txt(slide, sub, rx + Inches(0.35), ot + Inches(0.55),
            cw - Inches(0.55), Inches(0.45),
            size=Pt(13), color=LBLUE)
        ot += Inches(1.22)


def s12_recovery(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Recovery")

    txt(slide, "Strategy: re-approach and retry — up to MAX_RETRIES = 3",
        Inches(0.45), bar_h + Inches(0.3), W - Inches(0.9), Inches(0.42),
        size=Pt(17), bold=True, color=NAVY)

    cases = [
        (NAVY,  "Layer 1 fail (width / DI8)",
                "Open gripper → return to pick approach → re-close → recheck"),
        (BLUE,  "Layer 2 fail (post-lift width)",
                "Object slipped on lift → open → descend → retry full grip"),
        (RGBColor(0x1E,0x6E,0x45), "Layer 3 fail (CLIP)",
                "Visual not confirmed → open → return to approach → retry from top"),
        (AMBER, "Slip during transit",
                "CLIP re-verify → if lost, return to pick and retry"),
    ]

    ct = bar_h + Inches(0.85)
    cw = W - Inches(0.9)
    for col, title, desc in cases:
        rect(slide, Inches(0.45), ct, Inches(0.5), Inches(0.9), fill=col)
        rect(slide, Inches(0.95), ct, cw - Inches(0.5), Inches(0.9),
             fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.8))
        txt(slide, title, Inches(1.08), ct + Inches(0.05),
            cw - Inches(0.65), Inches(0.38),
            size=Pt(14), bold=True, color=DARK)
        txt(slide, desc, Inches(1.08), ct + Inches(0.44),
            cw - Inches(0.65), Inches(0.38),
            size=Pt(13), color=MED)
        ct += Inches(1.02)

    # WIP badge
    wip_t = ct + Inches(0.15)
    rect(slide, Inches(0.45), wip_t, W - Inches(0.9), Inches(0.6), fill=AMBER)
    txt(slide,
        "⚠  Detection and retry loop: implemented.  "
        "Advanced recovery (adapt pick pose, re-scan): Work in Progress.",
        Inches(0.65), wip_t + Inches(0.1), W - Inches(1.3), Inches(0.45),
        size=Pt(14), bold=True, color=WHITE)

    # Currently working
    cw_t = wip_t + Inches(0.75)
    rect(slide, Inches(0.45), cw_t, W - Inches(0.9), Inches(0.6), fill=GREEN)
    txt(slide,
        "✓  Currently working: failure type identified, gripper opened, robot returns to pick approach, retries up to 3×.",
        Inches(0.65), cw_t + Inches(0.1), W - Inches(1.3), Inches(0.45),
        size=Pt(14), bold=True, color=WHITE)


def s13_pipeline(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "End-to-End Autonomous Pipeline")

    steps = [
        ("1\nDetect",    NAVY,  "YOLO + depth"),
        ("2\nLocate",    NAVY,  "ArUco → robot XYZ"),
        ("3\nApproach",  NAVY,  "100 mm above pick"),
        ("4\nGrip",      NAVY,  "Descend + close RG2"),
        ("5\nCheck L1",  BLUE,  "Width + DI8"),
        ("6\nLift",      BLUE,  "+200 mm"),
        ("7\nCheck L2",  BLUE,  "Post-lift width"),
        ("8\nCLIP",      RGBColor(0x1E,0x6E,0x45), "≥75% Holding"),
        ("9\nTransfer",  NAVY,  "Slip monitor ON"),
        ("10\nPlace",    NAVY,  "Lower + release"),
    ]

    sw = (W - Inches(0.9)) / 5
    sh = Inches(2.4)
    st = bar_h + Inches(0.3)

    for i, (label, col, sub) in enumerate(steps):
        row = i // 5; c = i % 5
        l = Inches(0.45) + c * sw
        t = st + row * (sh + Inches(0.3))

        rect(slide, l, t, sw - Inches(0.1), Inches(0.9), fill=col)
        txt(slide, label, l, t + Inches(0.04), sw - Inches(0.1), Inches(0.85),
            size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, l, t + Inches(0.9), sw - Inches(0.1), sh - Inches(0.9),
             fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
        txt(slide, sub, l + Inches(0.08), t + Inches(1.0),
            sw - Inches(0.25), sh - Inches(1.1),
            size=Pt(13), color=DARK, align=PP_ALIGN.CENTER)

        # Arrow to the right (within row, not last of row)
        if c < 4:
            txt(slide, "→", l + sw - Inches(0.12), t + Inches(0.25),
                Inches(0.15), Inches(0.5),
                size=Pt(18), bold=True, color=MED)

    bottom_t = st + 2 * (sh + Inches(0.3))
    txt(slide, "✓  Fully autonomous loop — no operator input between picks.",
        Inches(0.45), bottom_t + Inches(0.1), W - Inches(0.9), Inches(0.42),
        size=Pt(15), bold=True, color=GREEN)


def s14_testing(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Evaluation — 90-Trial Protocol")

    txt(slide, "30 successful picks  ·  30 failed picks  ·  30 slip events",
        Inches(0.45), bar_h + Inches(0.3), W - Inches(0.9), Inches(0.42),
        size=Pt(16), bold=True, color=NAVY)

    rows = [
        ("True Positive Rate",    "TP / (TP + FN)", "≥ 85%", "Success correctly confirmed"),
        ("True Negative Rate",    "TN / (TN + FP)", "≥ 85%", "Failure correctly detected"),
        ("False Positive Rate",   "FP / (FP + TN)", "< 20%", "Says success when actually failed"),
        ("Overall Accuracy",      "(TP + TN) / 90", "≥ 85%", "All correct decisions"),
        ("Recovery Success Rate", "Recovered / Attempted", "≥ 70%", "Failures resolved by recovery"),
        ("Task Completion Rate",  "Placements / 90", "≥ 70%", "End-to-end performance"),
    ]

    cws = [Inches(3.1), Inches(2.6), Inches(1.4), Inches(5.5)]
    rh  = Inches(0.58)
    tt  = bar_h + Inches(0.85)

    # Header
    x = Inches(0.45)
    for cw, label in zip(cws, ["Metric", "Formula", "Target", "Interpretation"]):
        rect(slide, x, tt, cw - Inches(0.03), rh, fill=NAVY)
        txt(slide, label, x + Inches(0.1), tt + Inches(0.1),
            cw - Inches(0.2), rh - Inches(0.1),
            size=Pt(13), bold=True, color=WHITE)
        x += cw

    for i, row in enumerate(rows):
        t = tt + (i + 1) * rh
        bg_c = LIGHT if i % 2 == 0 else WHITE
        x = Inches(0.45)
        for j, (cw, cell) in enumerate(zip(cws, row)):
            c = GREEN if j == 2 else bg_c
            tc = DARK if j != 2 else WHITE
            rect(slide, x, t, cw - Inches(0.03), rh, fill=c,
                 line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
            txt(slide, cell, x + Inches(0.1), t + Inches(0.1),
                cw - Inches(0.2), rh - Inches(0.1),
                size=Pt(12), color=tc, bold=(j==2))
            x += cw


def s15_results(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Results")

    rect(slide, Inches(0.45), bar_h + Inches(0.3),
         W - Inches(0.9), H - bar_h - Inches(0.5), fill=WHITE,
         line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))

    txt(slide, "[ INSERT RESULTS TABLE AND CHARTS HERE ]",
        Inches(0.65), bar_h + Inches(1.0), W - Inches(1.3), Inches(1.0),
        size=Pt(22), bold=True, color=AMBER, align=PP_ALIGN.CENTER)

    multiline_txt(slide, [
        "Report: TPR, TNR, FPR, Overall Accuracy, Recovery Rate, Task Completion Rate",
        "",
        "Show: detection-only pipeline vs. detection + recovery pipeline comparison",
        "",
        "Include: screenshots from demo video — successful pick, detected failure, recovery sequence",
    ],
    Inches(0.65), bar_h + Inches(2.2), W - Inches(1.3), Inches(3.5),
    size=Pt(16), color=MED)


def s16_limitations(prs):
    slide = blank(prs)
    bg(slide)
    bar_h = topbar(slide, "Limitations & Future Work")

    lims = [
        "No collision avoidance — robot does not detect obstacles in workspace",
        "Calibration drift — hand-eye re-calibration needed if camera shifts",
        "CLIP accuracy not 100% — poor lighting / unusual crop angles cause errors",
        "No raw force readout — only binary DI8 contact flag available from UR10",
        "Recovery assumes object stays in the same position between retries",
    ]
    future = [
        "Adapt pick pose on retry — not just re-approach same point",
        "Re-detect object position before each retry attempt",
        "Replace URP gripper control with direct force-torque sensor",
        "Generalise beyond cubes & cylinders to arbitrary objects",
        "Add collision-aware motion planning",
    ]

    cw = (W - Inches(1.35)) / 2
    for col_i, (head, col, items) in enumerate([
        ("Current Limitations", AMBER, lims),
        ("Future Work",         BLUE,  future),
    ]):
        l = Inches(0.45) + col_i * (cw + Inches(0.45))
        t = bar_h + Inches(0.3)
        rect(slide, l, t, cw, Inches(0.52), fill=col)
        txt(slide, head, l + Inches(0.15), t + Inches(0.08),
            cw - Inches(0.3), Inches(0.38),
            size=Pt(15), bold=True, color=WHITE)
        rect(slide, l, t + Inches(0.52), cw, H - t - Inches(0.82),
             fill=WHITE, line_color=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(1))
        for j, item in enumerate(items):
            txt(slide, "• " + item,
                l + Inches(0.15), t + Inches(0.7) + j * Inches(0.95),
                cw - Inches(0.3), Inches(0.85),
                size=Pt(15), color=DARK)


def s17_conclusion(prs):
    slide = blank(prs)
    rect(slide, 0, 0, W, H, fill=NAVY)
    rect(slide, 0, Inches(5.3), W, Inches(0.07), fill=BLUE)

    txt(slide, "What Was Built",
        Inches(0.7), Inches(0.7), Inches(12), Inches(0.7),
        size=Pt(36), bold=True, color=WHITE)

    summary = [
        "3-layer grasp failure detection: width  ·  DI8 contact  ·  CLIP visual",
        "Slip detection with automatic CLIP re-verify during transit",
        "CLIP ViT-B/32 + Logistic Regression linear probe (trained on real gripper images)",
        "Fully autonomous end-to-end loop with retry recovery",
        "Evaluation targets: TPR ≥ 85%  ·  FPR < 20%  ·  Task completion ≥ 70%",
    ]
    for i, s in enumerate(summary):
        txt(slide, "✓  " + s,
            Inches(0.7), Inches(1.6) + i * Inches(0.82), Inches(12), Inches(0.72),
            size=Pt(17), color=RGBColor(0xD0,0xE8,0xF8))

    txt(slide,
        "Myu Wai Shin  ·  PDE3802 Final Year Project  ·  Middlesex University  ·  2025/26",
        Inches(0.7), Inches(6.8), Inches(12), Inches(0.45),
        size=Pt(13), color=RGBColor(0x70,0x90,0xB0))


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════
def main():
    prs = prs_new()

    s01_title(prs)
    s02_intro(prs)
    s03_aims(prs)
    divider(prs, "Hardware & Software", "What the system is built on")
    s04_hardware(prs)
    s05_software(prs)
    divider(prs, "Object Detection", "YOLO + OBB on OAK-D Lite")
    s06_yolo(prs)
    divider(prs, "Pose Estimation", "Eye-in-hand calibration → robot frame")
    s07_pose(prs)
    divider(prs, "Failure Detection Pipeline", "Gripper signals · Visual verification · Slip detection")
    s08_io(prs)
    s09_detection_arch(prs)
    s10_clip(prs)
    s11_slip(prs)
    s12_recovery(prs)
    s13_pipeline(prs)
    divider(prs, "Testing & Results", "90-trial evaluation protocol")
    s14_testing(prs)
    s15_results(prs)
    s16_limitations(prs)
    s17_conclusion(prs)

    out = "FYP_Presentation_v2.pptx"
    prs.save(out)
    print(f"\n  Saved: {out}  ({len(prs.slides)} slides)\n")

if __name__ == "__main__":
    main()
